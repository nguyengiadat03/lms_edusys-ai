from fastapi import FastAPI, File, UploadFile, HTTPException, Request, Form
from fastapi.responses import JSONResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from contextlib import asynccontextmanager
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import os
import tempfile
import time
import json
import logging
import re
import gc
import psutil
import threading
import concurrent.futures
from typing import Generator

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# OCR libs
logger.info("Checking library availability...")

# Initialize availability flags
EASYOCR_AVAILABLE = False
PYMUPDF_AVAILABLE = False
PPTX_AVAILABLE = False
DOCX_AVAILABLE = False
GEMINI_AVAILABLE = False
WHISPER_AVAILABLE = False

try:
    import easyocr
    EASYOCR_AVAILABLE = True
    logger.info("EasyOCR library available")
except (ImportError, MemoryError) as e:
    logger.warning(f"EasyOCR not available: {e}")
    EASYOCR_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
    logger.info("PyMuPDF library available")
except ImportError as e:
    logger.warning(f"PyMuPDF not available: {e}")
    PYMUPDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
    logger.info("python-pptx library available")
except ImportError as e:
    logger.warning(f"python-pptx not available: {e}")
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
    logger.info("python-docx library available")
except ImportError as e:
    logger.warning(f"python-docx not available: {e}")
    DOCX_AVAILABLE = False

try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
    logger.info("google-generativeai library available")
except ImportError as e:
    logger.warning(f"google-generativeai not available: {e}")
    GEMINI_AVAILABLE = False

try:
    import whisper
    logger.info("Whisper library available")
    WHISPER_AVAILABLE = True
except ImportError as e:
    logger.warning(f"Whisper not available: {e}")
    WHISPER_AVAILABLE = False

logger.info(f"Library check complete: EasyOCR={EASYOCR_AVAILABLE}, PyMuPDF={PYMUPDF_AVAILABLE}, PPTX={PPTX_AVAILABLE}, DOCX={DOCX_AVAILABLE}, Gemini={GEMINI_AVAILABLE}, Whisper={WHISPER_AVAILABLE}")

# Pydantic models for structured responses
class DocumentProcessingRequest(BaseModel):
    filePath: str
    fileType: str  # 'pdf' | 'docx' | 'ppt' | 'pptx' | 'xlsx' | 'image' | 'audio' | 'video'
    fileName: str

class DocumentAnalysisResult(BaseModel):
    extractedText: str
    summary: str
    segments: List[str]
    suggestedLevel: str
    suggestedTopic: str
    autoTags: List[str]
    metadata: Dict[str, Any]

@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup
    init_models()
    yield
    # Shutdown (if needed)

app = FastAPI(lifespan=lifespan)

# Allow CORS for frontend testing if needed
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Lazy loading models - only load when needed
easyocr_model = None
genai_model = None
whisper_model = None

def get_easyocr_model():
    """Lazy load EasyOCR model only when needed"""
    global easyocr_model
    if easyocr_model is None and EASYOCR_AVAILABLE:
        try:
            mem_usage = get_memory_usage()
            if mem_usage > 80:
                logger.warning(f"Memory too high ({mem_usage:.1f}%) for EasyOCR, skipping")
                return None
            logger.info("Loading EasyOCR model...")
            easyocr_model = easyocr.Reader(['en'])
            logger.info("EasyOCR model loaded successfully")
            force_gc()  # Clean up after loading
        except Exception as e:
            logger.error(f"Failed to load EasyOCR: {e}")
            easyocr_model = None
    return easyocr_model

def get_genai_model():
    """Lazy load Gemini model only when needed"""
    global genai_model
    if genai_model is None and GEMINI_AVAILABLE:
        try:
            # Try to load from .env file in current directory
            env_path = os.path.join(os.path.dirname(__file__), '.env')
            if os.path.exists(env_path):
                from dotenv import load_dotenv
                load_dotenv(env_path)
                logger.info("Loaded environment variables from .env file")

            api_key = os.getenv('GEMINI_API_KEY')
            if api_key:
                genai.configure(api_key=api_key)
                genai_model = genai.GenerativeModel('gemini-2.5-flash')
                logger.info("Gemini AI model loaded successfully")
            else:
                logger.warning("GEMINI_API_KEY not found")
                genai_model = None
        except Exception as e:
            logger.error(f"Failed to load Gemini AI: {e}")
            genai_model = None
    return genai_model

def get_whisper_model():
    """Lazy load Whisper model only when needed"""
    global whisper_model
    if whisper_model is None and WHISPER_AVAILABLE:
        try:
            # Check memory BEFORE loading
            mem_usage = get_memory_usage()
            if mem_usage > 75:  # Lower threshold for Whisper
                logger.warning(f"Memory too high ({mem_usage:.1f}%) for Whisper, skipping")
                return None

            # Check available RAM (need at least 2GB free)
            try:
                import psutil
                available_ram = psutil.virtual_memory().available / (1024**3)  # GB
                if available_ram < 2.0:
                    logger.warning(f"Insufficient RAM ({available_ram:.1f}GB) for Whisper")
                    return None
            except:
                pass  # Continue if psutil not available

            logger.info("Loading Whisper model...")
            # Add FFmpeg to PATH for Whisper
            ffmpeg_path = r"E:\edusys-ai3\backend\ffmpeg\ffmpeg-8.0-essentials_build\bin"
            current_path = os.environ.get('PATH', '')
            if ffmpeg_path not in current_path:
                os.environ['PATH'] = ffmpeg_path + ';' + current_path
                logger.info(f"Added FFmpeg to PATH: {ffmpeg_path}")

            import whisper
            whisper_model = whisper.load_model("base")
            logger.info("Whisper model loaded successfully")
            force_gc()  # Clean up after loading
        except Exception as e:
            logger.error(f"Failed to load Whisper: {e}")
            whisper_model = None
    return whisper_model

def unload_models():
    """Unload all models to free memory"""
    global easyocr_model, genai_model, whisper_model
    try:
        if easyocr_model is not None:
            del easyocr_model
            easyocr_model = None
            logger.info("EasyOCR model unloaded")
    except:
        pass

    try:
        if whisper_model is not None:
            del whisper_model
            whisper_model = None
            logger.info("Whisper model unloaded")
    except:
        pass

    # Keep Gemini loaded as it's lightweight
    force_gc()
    logger.info("Models unloaded, memory cleaned")

def init_models():
    """Minimal initialization - just check availability"""
    logger.info("Checking model availability (lazy loading enabled)")
    logger.info(f"Model availability: EasyOCR={EASYOCR_AVAILABLE}, Gemini={GEMINI_AVAILABLE}, Whisper={WHISPER_AVAILABLE}")

def clean_text(text):
    import re

    # Giữ lại dòng "--- Page x ---"
    def preserve_page_headers(match):
        return f"\n{match.group(0)}\n"

    # Bọc lại page headers bằng dòng riêng
    text = re.sub(r'--- Page \d+ ---', preserve_page_headers, text)

    # Xoá nhiều \n liên tục thành 1
    text = re.sub(r'\n+', '\n', text)

    # Gộp các dòng (trừ khi là page header)
    lines = text.split('\n')
    cleaned_lines = []
    for line in lines:
        if re.match(r'^--- Page \d+ ---$', line.strip()):
            cleaned_lines.append(line.strip())  # Giữ nguyên header
        elif line.strip():  # Nếu là dòng có nội dung
            cleaned_lines.append(line.strip())
    cleaned_text = ' '.join(cleaned_lines)

    # Xoá kí tự không in được
    cleaned_text = re.sub(r'[^\x20-\x7E]', '', cleaned_text)
    return cleaned_text.strip()


def prepare_text_for_gemini(text: str) -> str:
    """
    Hàm xử lý text đầu ra từ OCR/STT/PDF/... trước khi đưa qua Gemini AI
    """
    return clean_text(text)


def get_memory_usage() -> float:
    """Get current process memory usage percentage"""
    try:
        # Get current process memory usage
        process = psutil.Process()
        process_memory = process.memory_percent()
        return process_memory
    except:
        # Fallback to system memory if process memory fails
        try:
            return psutil.virtual_memory().percent
        except:
            return 0.0

def force_gc():
    """Force garbage collection with logging"""
    before = get_memory_usage()
    gc.collect()
    after = get_memory_usage()
    if before > after + 1:  # Only log if there's a meaningful reduction
        logger.info(f"GC: Memory {before:.1f}% -> {after:.1f}%")

def log_memory_status(message: str = ""):
    """Log current memory status"""
    mem = get_memory_usage()
    logger.info(f"Memory status {message}: {mem:.1f}%")

def chunk_text_stream(text: str, chunk_size: int = 4000) -> Generator[str, None, None]:
    """
    Chia text thành các chunks với memory-efficient streaming và aggressive cleanup
    """
    if len(text) <= chunk_size:
        yield text
        # Force garbage collection after yielding large text
        gc.collect()
        return

    words = text.split()
    current_chunk = ""
    word_count = 0

    for word in words:
        word_count += 1
        if len(current_chunk) + len(word) + 1 <= chunk_size:
            current_chunk += word + " "
        else:
            if current_chunk:
                chunk_to_yield = current_chunk.strip()
                current_chunk = ""  # Clear immediately
                yield chunk_to_yield
                # Force GC every 100 chunks to prevent memory buildup
                if word_count % 100 == 0:
                    gc.collect()
            current_chunk = word + " "

    if current_chunk:
        yield current_chunk.strip()

    # Final cleanup
    del words
    gc.collect()

def chunk_text(text: str, chunk_size: int = 4000) -> list:
    """
    Chia text thành các chunks với kích thước tối đa (backward compatibility)
    """
    return list(chunk_text_stream(text, chunk_size))


def analyze_chunk_with_gemini(chunk_text: str, chunk_index: int, tasks: list) -> dict:
    """
    Analyze một chunk với tất cả tasks trong một prompt duy nhất, memory-optimized với aggressive cleanup
    """
    if not GEMINI_AVAILABLE or genai_model is None:
        return {
            "chunk_index": chunk_index,
            "success": False,
            "error": "Gemini AI not available"
        }

    try:
        # Memory check before processing - stricter threshold
        mem_usage = get_memory_usage()
        if mem_usage > 85:  # Critical memory threshold
            force_gc()
            # Check again after GC
            mem_after_gc = get_memory_usage()
            if mem_after_gc > 90:
                logger.warning(f"Memory still too high after GC ({mem_after_gc:.1f}%) for chunk {chunk_index}, skipping")
                return {
                    "chunk_index": chunk_index,
                    "success": False,
                    "error": f"Memory usage too high ({mem_after_gc:.1f}%)"
                }
            logger.warning(f"High memory usage ({mem_usage:.1f}%) before chunk {chunk_index}, forced GC")

        # --- Build optimized prompt for all tasks ---
        task_instructions = []
        for task in tasks:
            t = task.lower()
            if t == "segment":
                task_instructions.append("SEGMENT: Create 3–5 short segment titles (3–8 words each) describing the main parts of the text, not a summary sentence. Example: ['Introduction paragraph', 'Main content about friendship', 'Description of activities', 'Conclusion about trust']")
            elif t == "level_suggestion":
                task_instructions.append("LEVEL: Suggest CEFR level (A1, A2, B1, B2, C1, C2)")
            elif t == "topic_suggestion":
                task_instructions.append("TOPIC: Suggest main topic or subject area")
            elif t == "tag_suggestion":
                task_instructions.append("TAGS: Suggest 4-5 relevant tags or keywords")
            elif t in ["summarize", "summary"]:
                task_instructions.append("SUMMARY: Provide 1-2 sentence summary")

        prompt = f"""Analyze this document and create a structured summary.

Document text:
{chunk_text}

Tasks to perform:
{chr(10).join(task_instructions)}

Return a JSON object with these exact fields:
{{
  "chunk_index": {chunk_index},
  "segments": ["segment title 1", "segment title 2", "segment title 3"],
  "level": "B1",
  "topic": "Main Topic",
  "suggested_tags": ["tag1", "tag2", "tag3"],
  "summary": "Brief summary text"
}}

CRITICAL REQUIREMENTS:
- segments must be an array of 3-5 SHORT titles (3-8 words each)
- level must be a CEFR level like A1, A2, B1, B2, C1, C2
- topic must be 1-3 words
- suggested_tags must be an array of strings
- summary must be 1-2 sentences

Return ONLY valid JSON, no markdown or extra text."""

        # --- Use threading with memory monitoring ---
        result = {"response": None, "error": None}

        def generate_with_timeout():
            try:
                result["response"] = genai_model.generate_content(prompt)
            except Exception as e:
                result["error"] = str(e)

        # Adjust timeout based on task complexity
        complex_tasks = ["segment", "level_suggestion", "topic_suggestion", "tag_suggestion"]
        has_complex_task = any(task.lower() in complex_tasks for task in tasks)
        timeout_seconds = 30.0 if has_complex_task else 15.0  # Longer timeout for complex tasks

        thread = threading.Thread(target=generate_with_timeout)
        thread.start()
        thread.join(timeout=timeout_seconds)

        # Force GC after threading operation
        force_gc()

        if thread.is_alive():
            return {
                "chunk_index": chunk_index,
                "success": False,
                "error": f"AI analysis timed out after {timeout_seconds} seconds"
            }

        if result["error"]:
            raise Exception(result["error"])

        response = result["response"]

        # --- Extract and parse response ---
        text_result = ""
        if hasattr(response, 'candidates') and response.candidates:
            candidate = response.candidates[0]
            if hasattr(candidate, 'content') and hasattr(candidate.content, 'parts'):
                for part in candidate.content.parts:
                    if hasattr(part, 'text'):
                        text_result += part.text

        if text_result.strip():
            logger.info(f"Gemini raw response for chunk {chunk_index}: '{text_result[:300]}...'")
            try:
                # Clean the response - remove markdown code blocks if present
                cleaned_response = text_result.strip()

                # Handle multiple possible markdown formats
                if cleaned_response.startswith('```json'):
                    cleaned_response = cleaned_response[7:]
                elif cleaned_response.startswith('```'):
                    cleaned_response = cleaned_response[3:]

                # Remove closing ```
                if cleaned_response.endswith('```'):
                    cleaned_response = cleaned_response[:-3]

                # Remove any remaining ``` markers
                cleaned_response = cleaned_response.replace('```', '')

                cleaned_response = cleaned_response.strip()

                # Try to parse the cleaned JSON
                parsed_json = json.loads(cleaned_response)
                parsed_json["chunk_index"] = chunk_index
                parsed_json["success"] = True

                # Clear large variables immediately
                del text_result, prompt, response, result, task_instructions
                force_gc()

                return parsed_json
            except json.JSONDecodeError as e:
                logger.error(f"JSON parse error for chunk {chunk_index}: {e}")
                logger.error(f"Raw response: {text_result[:500]}")

                # Try to extract individual fields from malformed response
                extracted_data = {
                    "chunk_index": chunk_index,
                    "success": True,
                    "segments": [],
                    "level": "",
                    "topic": "",
                    "suggested_tags": [],
                    "summary": ""
                }

                # Look for summary in the raw text
                if tasks and any(t.lower() in ["summarize", "summary"] for t in tasks):
                    # Try to find summary in the response
                    summary_match = re.search(r'"summary"\s*:\s*"([^"]*(?:\\"[^"]*)*[^"]*)"', text_result, re.IGNORECASE | re.DOTALL)
                    if summary_match:
                        extracted_data["summary"] = summary_match.group(1).replace('\\"', '"')
                    else:
                        # Use the whole response as summary if it's short enough
                        if len(text_result.strip()) < 500:
                            extracted_data["summary"] = text_result.strip()

                # Look for level
                if any(t.lower() == "level_suggestion" for t in tasks):
                    level_match = re.search(r'"level"\s*:\s*"([^"]*)"', text_result, re.IGNORECASE)
                    if level_match:
                        extracted_data["level"] = level_match.group(1)

                # Look for topic
                if any(t.lower() == "topic_suggestion" for t in tasks):
                    topic_match = re.search(r'"topic"\s*:\s*"([^"]*)"', text_result, re.IGNORECASE)
                    if topic_match:
                        extracted_data["topic"] = topic_match.group(1)

                # Look for segments
                if any(t.lower() == "segment" for t in tasks):
                    segments_match = re.search(r'"segments"\s*:\s*\[([^\]]*)\]', text_result, re.IGNORECASE | re.DOTALL)
                    if segments_match:
                        segments_text = segments_match.group(1)
                        # Extract quoted strings
                        segment_matches = re.findall(r'"([^"]*(?:\\"[^"]*)*[^"]*)"', segments_text)
                        extracted_data["segments"] = [s.replace('\\"', '"') for s in segment_matches]

                # Look for tags
                if any(t.lower() == "tag_suggestion" for t in tasks):
                    tags_match = re.search(r'"suggested_tags"\s*:\s*\[([^\]]*)\]', text_result, re.IGNORECASE | re.DOTALL)
                    if tags_match:
                        tags_text = tags_match.group(1)
                        # Try to parse tag objects
                        tag_objects = []
                        tag_matches = re.findall(r'\{\s*"tag_label"\s*:\s*"([^"]*(?:\\"[^"]*)*[^"]*)"\s*(?:,\s*"confidence"\s*:\s*([0-9.]+))?\s*\}', tags_text, re.IGNORECASE)
                        for tag_match in tag_matches:
                            confidence = float(tag_match[1]) if len(tag_match) > 1 and tag_match[1] else 0.8
                            tag_objects.append({
                                "tag_label": tag_match[0].replace('\\"', '"'),
                                "confidence": confidence
                            })
                        extracted_data["suggested_tags"] = tag_objects

                # If we extracted at least one field successfully, return the data
                if (extracted_data["summary"] or extracted_data["level"] or extracted_data["topic"] or
                    extracted_data["segments"] or extracted_data["suggested_tags"]):
                    return extracted_data

                # Fallback - return error
                return {
                    "chunk_index": chunk_index,
                    "success": False,
                    "error": "Failed to parse JSON response",
                    "raw_response": text_result[:200]  # Reduced size
                }
        else:
            return {
                "chunk_index": chunk_index,
                "success": False,
                "error": "Empty response from Gemini AI"
            }

    except Exception as e:
        logger.error(f"Gemini chunk analysis error: {e}")
        return {
            "chunk_index": chunk_index,
            "success": False,
            "error": f"AI analyze error: {str(e)}"
        }
    finally:
        # Ensure cleanup - clear all local variables
        try:
            del chunk_text, tasks
        except NameError:
            pass
        force_gc()


def analyze_with_gemini(text: str, task: str = "Extract and structure information from this document") -> dict:
    """
    Analyze extracted text with Gemini AI using chunking and multi-task approach
    """
    if not GEMINI_AVAILABLE or genai_model is None:
        return {
            "success": False,
            "error": "Gemini AI not available",
            "task": task,
            "content": text[:500] + "..." if len(text) > 500 else text
        }

    try:
        start_time = time.time()

        # Convert single task to list for compatibility
        tasks = [task] if isinstance(task, str) else task

        # For short text, use single chunk with memory protection
        if len(text) <= 4000:
            # Check memory before single chunk processing
            mem_usage = get_memory_usage()
            if mem_usage > 80:
                force_gc()
                mem_after = get_memory_usage()
                if mem_after > 85:
                    logger.warning(f"Memory too high for single chunk: {mem_after:.1f}%, attempting summary with truncated text")
                    truncated_text = text[:500] + "..." if len(text) > 500 else text
                    # Attempt summary with truncated text
                    try:
                        summary_result = analyze_chunk_with_gemini(truncated_text, 0, ["summarize"])
                        if summary_result.get("success"):
                            return {
                                "success": True,
                                "summary": summary_result.get("summary", ""),
                                "processing_time": f"{time.time() - start_time:.2f}s",
                                "note": "Summary generated from truncated text due to high memory usage"
                            }
                    except Exception as e:
                        logger.error(f"Failed to generate summary from truncated text: {e}")

                    # Fallback if summary fails
                    return {
                        "success": False,
                        "error": f"Memory usage too high ({mem_after:.1f}%) for processing",
                        "processing_time": f"{time.time() - start_time:.2f}s",
                        "fallback_text": truncated_text
                    }

            # Check if tasks include complex analysis that might need more time
            complex_tasks = ["segment", "level_suggestion", "topic_suggestion", "tag_suggestion"]
            has_complex_task = any(task.lower() in complex_tasks for task in tasks)
            if has_complex_task and len(text) > 2000:
                # For complex tasks with longer text, try to process in smaller chunks
                logger.info(f"Complex task detected, processing text in smaller chunks")
                chunks = chunk_text(text, 2000)  # Smaller chunks for complex tasks
                if len(chunks) > 1:
                    chunk_results = []
                    for i, chunk in enumerate(chunks[:3]):  # Limit to first 3 chunks
                        try:
                            chunk_result = analyze_chunk_with_gemini(chunk, i, tasks)
                            chunk_results.append(chunk_result)
                        except Exception as e:
                            logger.error(f"Failed to process chunk {i}: {e}")
                            break

                    if chunk_results:
                        # Merge results from chunks
                        merged_result = {
                            "success": True,
                            "processing_time": f"{time.time() - start_time:.2f}s",
                            "chunks_processed": len(chunk_results)
                        }

                        # Merge segments, levels, topics, tags, summaries
                        all_segments = []
                        all_tags = []
                        summaries = []
                        levels = []
                        topics = []

                        for result in chunk_results:
                            if result.get("success"):
                                if "segments" in result and result["segments"]:
                                    all_segments.extend(result["segments"])
                                if "suggested_tags" in result and result["suggested_tags"]:
                                    all_tags.extend(result["suggested_tags"])
                                if "summary" in result and result["summary"]:
                                    summaries.append(result["summary"])
                                if "level" in result and result["level"]:
                                    levels.append(result["level"])
                                if "topic" in result and result["topic"]:
                                    topics.append(result["topic"])

                        if all_segments:
                            merged_result["segments"] = list(set(all_segments))[:10]  # Limit segments
                        if all_tags:
                            # Group and average confidence for tags
                            tag_groups = {}
                            for tag in all_tags:
                                label = tag.get("tag_label", "")
                                conf = tag.get("confidence", 0.5)
                                if label in tag_groups:
                                    tag_groups[label]["confidence"] = (tag_groups[label]["confidence"] + conf) / 2
                                    tag_groups[label]["count"] += 1
                                else:
                                    tag_groups[label] = {"confidence": conf, "count": 1}
                            merged_result["suggested_tags"] = [
                                {"tag_label": label, "confidence": data["confidence"]}
                                for label, data in sorted(tag_groups.items(), key=lambda x: x[1]["confidence"], reverse=True)[:5]
                            ]
                        if summaries:
                            merged_result["summary"] = " ".join(summaries)
                        if levels:
                            merged_result["level"] = max(set(levels), key=levels.count)
                        if topics:
                            merged_result["topic"] = max(set(topics), key=topics.count)

                        return merged_result

            result = analyze_chunk_with_gemini(text, 0, tasks)
            result["processing_time"] = f"{time.time() - start_time:.2f}s"
            return result

        # Enable chunking for large texts - memory safe with proper chunking
        # No artificial truncation needed since chunking handles memory management

        # Use chunking with larger chunks for better performance
        logger.info(f"Processing text ({len(text)} chars) with optimized chunking")

        # Use larger chunks (8000 chars) to reduce API calls while staying within limits
        chunks = chunk_text(text, 8000)
        logger.info(f"Split into {len(chunks)} chunks (8000 chars each)")

        chunk_results = []
        # Process first 3 chunks for good coverage without too many API calls
        for i, chunk in enumerate(chunks[:3]):
            try:
                chunk_result = analyze_chunk_with_gemini(chunk, i, tasks)
                chunk_results.append(chunk_result)
            except Exception as e:
                logger.error(f"Failed to process chunk {i}: {e}")
                break

        if chunk_results:
            # Merge results intelligently
            merged_result = {
                "success": True,
                "processing_time": f"{time.time() - start_time:.2f}s",
                "chunks_processed": len(chunk_results)
            }

            # Collect all results
            all_segments = []
            all_tags = []
            summaries = []
            levels = []
            topics = []

            for result in chunk_results:
                if result.get("success"):
                    if "segments" in result and result["segments"]:
                        all_segments.extend(result["segments"])
                    if "suggested_tags" in result and result["suggested_tags"]:
                        all_tags.extend(result["suggested_tags"])
                    if "summary" in result and result["summary"]:
                        summaries.append(result["summary"])
                    if "level" in result and result["level"]:
                        levels.append(result["level"])
                    if "topic" in result and result["topic"]:
                        topics.append(result["topic"])

            # Merge segments (remove duplicates, limit to 5)
            if all_segments:
                merged_result["segments"] = list(set(all_segments))[:5]

            # Merge tags (group by tag, average confidence, top 5)
            if all_tags:
                tag_groups = {}
                for tag in all_tags:
                    if isinstance(tag, dict) and "tag_label" in tag:
                        label = tag["tag_label"]
                        conf = tag.get("confidence", 0.8)
                    else:
                        label = str(tag)
                        conf = 0.8

                    if label in tag_groups:
                        tag_groups[label]["confidence"] = (tag_groups[label]["confidence"] + conf) / 2
                        tag_groups[label]["count"] += 1
                    else:
                        tag_groups[label] = {"confidence": conf, "count": 1}

                merged_result["suggested_tags"] = [
                    {"tag_label": label, "confidence": data["confidence"]}
                    for label, data in sorted(tag_groups.items(), key=lambda x: x[1]["confidence"], reverse=True)[:5]
                ]

            # Combine summaries
            if summaries:
                merged_result["summary"] = " ".join(summaries)

            # Use most common level/topic
            if levels:
                merged_result["level"] = max(set(levels), key=levels.count)
            if topics:
                merged_result["topic"] = max(set(topics), key=topics.count)

            return merged_result

        return {
            "success": False,
            "error": "No chunks could be processed",
            "processing_time": f"{time.time() - start_time:.2f}s"
        }

        # For small texts, use single chunk
        result = analyze_chunk_with_gemini(text, 0, tasks)
        result["processing_time"] = f"{time.time() - start_time:.2f}s"

        # Clear text variable immediately after processing
        del text
        force_gc()

        return result

        # DISABLED: Parallel chunking for now due to memory issues
        """
        # For longer text, use chunking with controlled parallelism
        chunks = chunk_text(text, 4000)
        max_concurrent = min(3, len(chunks))  # Max 3 concurrent chunks

        import concurrent.futures
        chunk_results = []

        # Process chunks in batches to control memory usage
        for i in range(0, len(chunks), max_concurrent):
            batch_chunks = chunks[i:i + max_concurrent]

            with concurrent.futures.ThreadPoolExecutor(max_workers=max_concurrent) as executor:
                futures = [
                    executor.submit(analyze_chunk_with_gemini, chunk, chunk_index + i, tasks)
                    for chunk_index, chunk in enumerate(batch_chunks)
                ]

                for future in concurrent.futures.as_completed(futures, timeout=60.0):
                    try:
                        chunk_results.append(future.result())
                    except concurrent.futures.TimeoutError:
                        chunk_results.append({
                            "chunk_index": -1,
                            "success": False,
                            "error": "Chunk processing timeout"
                        })
        """

        # Merge results from all chunks
        merged_result = {
            "success": True,
            "processing_time": f"{time.time() - start_time:.2f}s",
            "total_chunks": len(chunks),
            "successful_chunks": sum(1 for r in chunk_results if r.get("success", False))
        }

        # Merge segments from all chunks
        all_segments = []
        for result in chunk_results:
            if result.get("success") and "segments" in result:
                all_segments.extend(result["segments"])
        if all_segments:
            merged_result["segments"] = list(set(all_segments))  # Remove duplicates

        # Use the most common level/topic from chunks
        levels = [r.get("level") for r in chunk_results if r.get("success") and "level" in r]
        if levels:
            merged_result["level"] = max(set(levels), key=levels.count)

        topics = [r.get("topic") for r in chunk_results if r.get("success") and "topic" in r]
        if topics:
            merged_result["topic"] = max(set(topics), key=topics.count)

        # Merge and deduplicate tags
        all_tags = []
        for result in chunk_results:
            if result.get("success") and "suggested_tags" in result:
                all_tags.extend(result["suggested_tags"])

        if all_tags:
            # Group by tag_label and average confidence
            tag_groups = {}
            for tag in all_tags:
                label = tag.get("tag_label", "")
                conf = tag.get("confidence", 0.5)
                if label in tag_groups:
                    tag_groups[label]["confidence"] = (tag_groups[label]["confidence"] + conf) / 2
                    tag_groups[label]["count"] += 1
                else:
                    tag_groups[label] = {"confidence": conf, "count": 1}

            merged_result["suggested_tags"] = [
                {"tag_label": label, "confidence": data["confidence"]}
                for label, data in sorted(tag_groups.items(), key=lambda x: x[1]["confidence"], reverse=True)[:5]
            ]

        # Combine summaries
        summaries = [r.get("summary") for r in chunk_results if r.get("success") and "summary" in r]
        if summaries:
            merged_result["summary"] = " ".join(summaries)

        return merged_result

    except Exception as e:
        logger.error(f"Gemini AI analyze error: {e}")
        return {
            "success": False,
            "error": f"AI analyze error: {str(e)}",
            "task": task,
            "content": text[:500] + "..." if len(text) > 500 else text
        }


@app.get("/health")
async def health():
    return {
        'status': 'healthy',
        'services': {
            'easyocr': EASYOCR_AVAILABLE,
            'pymupdf': PYMUPDF_AVAILABLE,
            'pptx': PPTX_AVAILABLE,
            'docx': DOCX_AVAILABLE,
            'gemini': GEMINI_AVAILABLE,
            'whisper': WHISPER_AVAILABLE
        },
        'loaded_models': {
            'easyocr': easyocr_model is not None,
            'gemini': genai_model is not None,
            'whisper': whisper_model is not None
        },
        'memory_usage': f"{get_memory_usage():.1f}%"
    }

@app.post("/ocr/image")
async def ocr_image(file: UploadFile = File(...)):
    if not file or file.filename == '':
        raise HTTPException(400, "No file provided or selected")

    suffix = os.path.splitext(file.filename)[1]
    temp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_path = temp_file.name
            content = await file.read()
            temp_file.write(content)

        logger.info(f"OCR processing file: {file.filename}, size: {len(content)} bytes")

        extracted_text = ""

        if easyocr_model:
            try:
                results = easyocr_model.readtext(temp_path)
                logger.info(f"EasyOCR results: type={type(results)}, len={len(results) if results else 0}")
                if results and isinstance(results, list):
                    extracted_text = ' '.join([text for (_, text, _) in results])
                extracted_text = extracted_text.strip()

            except Exception as e:
                logger.error(f"EasyOCR error: {e}")

        words = extracted_text.split()
        if not extracted_text or len(words) < 1:
            return {
                'success': True,
                'extracted_text': '',
                'confidence': 'low',
                'method': 'none'
            }

        return {
            'success': True,
            'extracted_text': extracted_text,
            'confidence': 'medium',
            'method': 'easyocr'
        }

    except Exception as e:
        logger.error(f"OCR error: {e}")
        raise HTTPException(500, str(e))
    finally:
        if temp_path and os.path.exists(temp_path):
            os.unlink(temp_path)





@app.post("/stt/audio")
async def stt_audio(file: UploadFile = File(...)):
    if not file or file.filename == '':
        raise HTTPException(400, "No file provided or selected")

    suffix = os.path.splitext(file.filename)[1]
    temp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_path = temp_file.name
            content = await file.read()
            temp_file.write(content)

        logger.info(f"Saved audio file {file.filename} at {temp_path}")

        if not whisper_model:
            raise HTTPException(503, "Whisper model not available")

        time.sleep(0.1)  # ensure file fully written

        # Check memory before processing
        mem_before = get_memory_usage()
        if mem_before > 80:
            logger.warning(f"High memory before Whisper processing: {mem_before:.1f}%")
            force_gc()
            mem_after_gc = get_memory_usage()
            if mem_after_gc > 85:
                raise HTTPException(503, f"Memory too high for audio processing: {mem_after_gc:.1f}%")

        try:
            result = whisper_model.transcribe(temp_path, fp16=False)
            transcription = result.get('text', '').strip()
            language = result.get('language', 'unknown')

            # Force cleanup immediately after processing
            del result
            force_gc()

            # Check memory after processing and unload if needed
            mem_after = get_memory_usage()
            if mem_after > 80:
                logger.warning(f"High memory after Whisper processing: {mem_after:.1f}%, unloading model")
                unload_models()

        except Exception as transcribe_error:
            # Force cleanup on error
            force_gc()
            raise transcribe_error

        logger.info(f"Whisper transcription done, length={len(transcription)}, text: '{transcription[:100]}...'")

        # Check if transcription is just metadata/filename (common Whisper fallback)
        # Be less strict - allow short transcriptions that might be valid
        if not transcription or transcription.strip() == '':
            logger.warning(f"Whisper returned empty transcription, treating as failed")
            raise HTTPException(503, "Whisper transcription failed - returned empty text")
        elif len(transcription.strip()) < 2:
            logger.warning(f"Whisper returned very short transcription: '{transcription}', might be poor quality")
            # Don't fail for short transcriptions, just log warning
        elif transcription.lower().strip() in [file.filename.lower(), 'track', 'audio', 'sound', 'music', 'song']:
            logger.warning(f"Whisper returned generic audio term: '{transcription}', likely poor quality but allowing")
            # Allow generic terms but log warning

        return {
            'success': True,
            'transcription': transcription,
            'language': language,
            'confidence': 'high',
            'segments': result.get('segments', [])
        }

    except Exception as e:
        logger.error(f"STT error: {e}")
        raise HTTPException(500, str(e))
    finally:
        if temp_path and os.path.exists(temp_path):
            os.unlink(temp_path)

@app.post("/extract/pdf")
async def extract_pdf(file: UploadFile = File(...)):
    if not file or file.filename == '':
        raise HTTPException(400, "No file provided or selected")

    if not PYMUPDF_AVAILABLE:
        raise HTTPException(503, "PyMuPDF not available")

    temp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
            temp_path = temp_file.name
            content = await file.read()
            temp_file.write(content)

        import fitz
        doc = fitz.open(temp_path)
        extracted_text = ""
        page_count = len(doc)

        for page_num in range(page_count):
            try:
                page = doc.load_page(page_num)
                text = page.get_text()
                extracted_text += text + " "
            except Exception as e:
                logger.warning(f"Failed to extract text from page {page_num + 1}: {e}")
                extracted_text += "[Error extracting text] "

        doc.close()

        cleaned = clean_text(extracted_text)

        if not cleaned or cleaned == "[Error extracting text]":
            return JSONResponse(status_code=422, content={
                'success': False,
                'error': 'No text could be extracted from PDF',
                'extracted_text': '',
                'pages': page_count,
                'confidence': 'none'
    })

        return {
               'success': True,
               'extracted_text': cleaned,
               'pages': page_count,
               'confidence': 'high' if len(cleaned.split()) > 10 else 'medium'
}

    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        raise HTTPException(500, str(e))
    finally:
        if temp_path and os.path.exists(temp_path):
            os.unlink(temp_path)

@app.post("/extract/ppt")
async def extract_ppt(file: UploadFile = File(...)):
    if not file or file.filename == '':
        raise HTTPException(400, "No file provided or selected")

    if not PPTX_AVAILABLE:
        raise HTTPException(503, "python-pptx not available")

    suffix = '.pptx' if file.filename.lower().endswith('.pptx') else '.ppt'
    temp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_path = temp_file.name
            content = await file.read()
            temp_file.write(content)

        presentation = Presentation(temp_path)
        extracted_text = ""
        slide_count = 0

        for idx, slide in enumerate(presentation.slides):
            slide_count += 1
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + " "

        return {
            'success': True,
            'extracted_text': extracted_text.strip(),
            'slides': slide_count,
            'confidence': 'high' if extracted_text.strip() else 'low'
        }

    except Exception as e:
        logger.error(f"PPT extraction error: {e}")
        raise HTTPException(500, str(e))
    finally:
                if temp_path and os.path.exists(temp_path):
                  os.unlink(temp_path)

@app.post("/ai/analyze")
async def ai_analyze(request: Request):
    """
    Nhận JSON gồm 'task' và 'content' để gọi Gemini AI phân tích
    """
    if not GEMINI_AVAILABLE or genai_model is None:
        raise HTTPException(503, "Gemini AI model not available")

    try:
        data = await request.json()
        task = data.get('task')
        content = data.get('content')

        if not task or not content:
            raise HTTPException(400, "Missing 'task' or 'content' in request body")

        prompt = f"{task}\n\n{content}"

        response = genai_model.generate_content(prompt)

        if response.text:
            text_result = response.text
        else:
            text_result = ""

        return {
            'success': True,
            'task': task,
            'content': content,
            'analysis': text_result
        }

    except Exception as e:
        logger.error(f"Gemini AI analyze error: {e}")
        raise HTTPException(500, f"AI analyze error: {e}")

@app.post("/process/file")
async def process_file(file: UploadFile = File(...)):
    logger.info("process_file called - unified OCR + AI analysis")
    """
    Unified endpoint: Extract text from file -> Analyze with Gemini AI -> Return structured JSON with all tasks
    """
    task = "segment,level_suggestion,topic_suggestion,tag_suggestion,summary"  # Always do all tasks
    if not file or file.filename == '':
        raise HTTPException(400, "No file provided or selected")

    suffix = os.path.splitext(file.filename)[1].lower()
    temp_path = None
    start_time = time.time()

    try:
        content = await file.read()
        file_size_mb = len(content) / (1024 * 1024)

        # File size limits - INCREASED to support larger documents with chunking
        MAX_FILE_SIZE_MB = 50  # Increased to 50MB to support larger documents
        if file_size_mb > MAX_FILE_SIZE_MB:
            raise HTTPException(413, f"File too large: {file_size_mb:.1f}MB (max {MAX_FILE_SIZE_MB}MB)")

        logger.info(f"Processing file: {file.filename}, size: {file_size_mb:.1f}MB, type: {suffix}")

        extracted_text = ""

        # Extract text based on file type with lazy model loading
        if suffix in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.gif', '.webp']:
            # Image processing with lazy-loaded EasyOCR
            easyocr_model = get_easyocr_model()
            if not easyocr_model:
                return analyze_with_gemini("", f"Error: EasyOCR model not available for image {file.filename}")

            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                temp_path = temp_file.name
                temp_file.write(content)

            try:
                results = easyocr_model.readtext(temp_path, detail=0)
                logger.info(f"EasyOCR results: type={type(results)}, len={len(results) if results else 0}")
                if results and isinstance(results, list):
                    extracted_text = ' '.join([str(text) for text in results if text])
                extracted_text = extracted_text.strip()

            except Exception as ocr_error:
                logger.error(f"EasyOCR processing failed: {ocr_error}")
                return {
                    "summary": "",
                    "analysis_type": task.lower(),
                    "processing_time": f"{time.time() - start_time:.2f}s",
                    "error": f"OCR processing failed: {str(ocr_error)}"
                }

            # Check if OCR extracted meaningful text
            if not extracted_text or len(extracted_text.strip()) < 5:
                logger.warning(f"EasyOCR failed to extract text from image {file.filename}")
                return {
                    "summary": "",
                    "analysis_type": task.lower(),
                    "processing_time": f"{time.time() - start_time:.2f}s",
                    "error": "OCR could not extract readable text from the image"
                }

            # Limit extracted text length for performance
            MAX_IMAGE_TEXT = 3000
            if len(extracted_text) > MAX_IMAGE_TEXT:
                extracted_text = extracted_text[:MAX_IMAGE_TEXT] + "... [truncated for performance]"

        elif suffix == '.pdf':
            # PDF processing with strict memory controls
            if not PYMUPDF_AVAILABLE:
                return analyze_with_gemini("", f"Error: PyMuPDF not available for PDF {file.filename}")

            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                temp_path = temp_file.name
                temp_file.write(content)

            import fitz
            doc = fitz.open(temp_path)
            page_count = len(doc)

            # Allow full PDF extraction since chunking handles memory management
            for page_num in range(page_count):
                try:
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    extracted_text += f"\n--- Page {page_num + 1} ---\n{text}\n"
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num + 1}: {e}")
                    extracted_text += f"\n--- Page {page_num + 1} ---\n[Error extracting text]\n"

            doc.close()

        elif suffix in ['.ppt', '.pptx']:
            # PPT processing with ultra-strict limits
            if not PPTX_AVAILABLE:
                return analyze_with_gemini("", f"Error: python-pptx not available for PPT {file.filename}")

            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                temp_path = temp_file.name
                temp_file.write(content)

            presentation = Presentation(temp_path)

            for idx, slide in enumerate(presentation.slides):
                slide_text = f"\n--- Slide {idx + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                extracted_text += slide_text

        elif suffix in ['.mp3', '.wav', '.mp4', '.avi', '.m4a', '.flac', '.ogg']:
            # Audio/Video processing with lazy-loaded Whisper
            whisper_model = get_whisper_model()
            if not whisper_model:
                return analyze_with_gemini("", f"Error: Whisper not available for audio/video {file.filename}")

            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                temp_path = temp_file.name
                temp_file.write(content)

            try:
                # Ensure file is fully written before transcription
                time.sleep(1.0)

                # Verify file exists and has content
                if not os.path.exists(temp_path):
                    raise Exception(f"Temp file not found: {temp_path}")

                file_size = os.path.getsize(temp_path)
                if file_size == 0:
                    raise Exception(f"Temp file is empty: {temp_path}")

                logger.info(f"Starting Whisper transcription for {file.filename} ({file_size} bytes)")

                result = whisper_model.transcribe(temp_path, fp16=False)
                transcription = result.get('text', '').strip()
                language = result.get('language', 'unknown')

                if not transcription or len(transcription.strip()) < 5:
                    logger.warning(f"Whisper returned poor transcription for {file.filename}")
                    return {
                        "summary": "",
                        "analysis_type": "audio_transcription",
                        "processing_time": f"{time.time() - start_time:.2f}s",
                        "error": "Audio transcription failed - poor quality or empty result"
                    }

                logger.info(f"Whisper transcription successful: {len(transcription)} chars, language: {language}")
                extracted_text = transcription

            except Exception as e:
                logger.error(f"Audio/video processing failed: {e}")
                return {
                    "summary": "",
                    "analysis_type": "audio_transcription",
                    "processing_time": f"{time.time() - start_time:.2f}s",
                    "error": f"Audio/video processing failed: {str(e)}"
                }

        elif suffix == '.docx':
            # DOCX processing with ultra-strict limits
            if not DOCX_AVAILABLE:
                return analyze_with_gemini("", f"Error: python-docx not available for DOCX {file.filename}")

            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                temp_path = temp_file.name
                temp_file.write(content)

            doc = Document(temp_path)

            for para in doc.paragraphs:
                extracted_text += para.text + "\n"

        else:
            return analyze_with_gemini("", f"Error: Unsupported file type {suffix} for {file.filename}")

        # Clean extracted text
        cleaned_text = clean_text(extracted_text)

        # Clear extracted_text immediately to free memory
        del extracted_text
        force_gc()

        # Log processing time for monitoring
        extraction_time = time.time() - start_time
        logger.info(f"Text extraction completed in {extraction_time:.2f}s, text length: {len(cleaned_text)}")
        logger.info(f"Extracted text preview: '{cleaned_text[:200]}...'")

        # Skip analysis if no text extracted
        if not cleaned_text or len(cleaned_text.strip()) < 10:
            logger.warning(f"No meaningful text extracted from {file.filename}")
            return {
                "summary": "",
                "analysis_type": task.lower(),
                "processing_time": f"{time.time() - start_time:.2f}s",
                "error": "No text could be extracted from the file"
            }

        # Analyze with lazy-loaded Gemini AI
        genai_model = get_genai_model()
        if not genai_model:
            return {
                "chunk_index": 0,
                "segments": [],
                "level": "",
                "topic": "",
                "suggested_tags": [],
                "summary": cleaned_text[:500] + "..." if len(cleaned_text) > 500 else cleaned_text,
                "success": True,
                "processing_time": f"{time.time() - start_time:.2f}s",
                "note": "Gemini AI not available, returning raw text as summary"
            }

        logger.info("Starting Gemini analysis for all tasks")
        result = analyze_with_gemini(cleaned_text, task)
        logger.info(f"Gemini analysis completed: success={result.get('success', False)}, result_keys={list(result.keys()) if isinstance(result, dict) else 'not_dict'}")

        # Ensure consistent output format for all file types
        if result.get("success"):
            # Standardize the output format
            standardized_result = {
                "chunk_index": result.get("chunk_index", 0),
                "segments": result.get("segments", []),
                "level": result.get("level", ""),
                "topic": result.get("topic", ""),
                "suggested_tags": result.get("suggested_tags", []),
                "summary": result.get("summary", ""),
                "success": True,
                "processing_time": result.get("processing_time", f"{time.time() - start_time:.2f}s")
            }
            # Add any additional fields from the original result
            for key, value in result.items():
                if key not in standardized_result:
                    standardized_result[key] = value
        else:
            # Return standardized error format
            standardized_result = {
                "chunk_index": 0,
                "segments": [],
                "level": "",
                "topic": "",
                "suggested_tags": [],
                "summary": result.get("summary", ""),
                "success": False,
                "error": result.get("error", "Analysis failed"),
                "processing_time": result.get("processing_time", f"{time.time() - start_time:.2f}s")
            }

        # Clear cleaned_text immediately after analysis to free memory
        del cleaned_text
        force_gc()

        # Unload heavy models after processing to free memory for next request
        unload_models()

        # Return the standardized result
        return standardized_result

    except Exception as e:
        logger.error(f"Processing error: {e}")
        return analyze_with_gemini("", f"Error processing file {file.filename}: {str(e)}")
    finally:
        if temp_path and os.path.exists(temp_path):
            os.unlink(temp_path)

@app.get("/")
async def main_page():
    html_content = """
    <!DOCTYPE html>
    <html>
    <head>
        <title>OCR + AI Service</title>
        <style>
            body { font-family: Arial, sans-serif; margin: 20px; background: #f5f5f5; }
            .container { max-width: 800px; margin: 0 auto; background: white; padding: 30px; border-radius: 10px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }
            .form-section { margin-bottom: 30px; padding: 20px; border: 1px solid #e0e0e0; border-radius: 8px; background: #fafafa; }
            .result { margin-top: 20px; padding: 15px; background: #f8f9fa; border: 1px solid #dee2e6; border-radius: 5px; white-space: pre-wrap; max-height: 500px; overflow-y: auto; font-family: monospace; font-size: 12px; }
            .status { padding: 5px 10px; border-radius: 3px; font-weight: bold; display: inline-block; margin: 5px; }
            .status.available { background: #d4edda; color: #155724; }
            .status.unavailable { background: #f8d7da; color: #721c24; }
            input[type="file"] { margin: 10px 0; padding: 8px; border: 1px solid #ccc; border-radius: 4px; width: 100%; }
            input[type="text"], textarea { margin: 10px 0; padding: 8px; border: 1px solid #ccc; border-radius: 4px; width: 100%; }
            button { background: #007bff; color: white; padding: 10px 20px; border: none; border-radius: 5px; cursor: pointer; margin: 5px; }
            button:hover { background: #0056b3; }
            .task-selector { margin: 15px 0; }
            .task-selector select { padding: 8px; border: 1px solid #ccc; border-radius: 4px; width: 100%; }
        </style>
        <script>
            async function checkHealth() {
                try {
                    const response = await fetch('/health');
                    const data = await response.json();
                    document.getElementById('health-result').innerHTML = '<pre>' + JSON.stringify(data, null, 2) + '</pre>';
                } catch (error) {
                    document.getElementById('health-result').innerHTML = 'Error: ' + error.message;
                }
            }

            async function submitUnifiedForm(form, resultDiv) {
                const formData = new FormData(form);

                resultDiv.innerHTML = '⏳ Processing...';

                try {
                    const response = await fetch('/process/file', {
                        method: 'POST',
                        body: formData
                    });
                    const data = await response.json();
                    resultDiv.innerHTML = '<pre>' + JSON.stringify(data, null, 2) + '</pre>';
                } catch (error) {
                    resultDiv.innerHTML = '❌ Error: ' + error.message;
                }
                return false;
            }

            // No task selection needed - always does full analysis

            window.onload = function() {
                checkHealth();
            };
        </script>
    </head>
    <body>
        <div class="container">
            <h1>📄 OCR + AI Document Processor</h1>
            <p>Upload any document (PDF, images, PPT, DOCX) and get AI-powered analysis instantly.</p>

            <div class="form-section">
                <h2>🔧 Service Status</h2>
                <button onclick="checkHealth()">🔄 Check Services</button>
                <div id="health-result" class="result"></div>
            </div>

            <div class="form-section">
                <h2>📤 Upload & Process Document</h2>
                <p>Supported formats: PDF, Images (JPG/PNG), PowerPoint (PPT/PPTX), Word (DOCX), Audio (MP3/WAV), Video (MP4/AVI)</p>

                <p><strong>AI Analysis:</strong> Automatic document segmentation, level assessment, topic identification, and tagging</p>

                <form action="/process/file" onsubmit="return submitUnifiedForm(this, document.getElementById('process-result'))" enctype="multipart/form-data" method="post">
                    <input name="file" type="file" accept=".pdf,.jpg,.jpeg,.png,.bmp,.tiff,.tif,.gif,.webp,.ppt,.pptx,.docx,.mp3,.wav,.mp4,.avi,.m4a,.flac,.ogg" required>
                    <br><br>
                    <button type="submit">🚀 Process Document</button>
                </form>
                <div id="process-result" class="result"></div>
            </div>

            <div class="form-section">
                <h2>ℹ️ How it works</h2>
                <ol>
                    <li><strong>Upload:</strong> Select any supported document file</li>
                    <li><strong>OCR:</strong> Text is automatically extracted from the file</li>
                    <li><strong>AI Analysis:</strong> Content is analyzed by Gemini AI</li>
                    <li><strong>Result:</strong> Get structured JSON with insights and summary</li>
                </ol>
                <p><strong>Supported file types:</strong> PDF, JPG, PNG, BMP, TIFF, GIF, WEBP, PPT, PPTX, DOCX, MP3, WAV, MP4, AVI, M4A, FLAC, OGG</p>
            </div>
        </div>
    </body>
    </html>
    """
    return HTMLResponse(content=html_content, status_code=200)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="localhost", port=8000)

